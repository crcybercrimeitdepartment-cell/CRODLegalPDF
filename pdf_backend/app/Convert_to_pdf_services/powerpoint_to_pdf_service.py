"""
PowerPoint to PDF conversion service.
Uses python-pptx + reportlab — no Microsoft PowerPoint installation required.
"""

import io
import logging
from pathlib import Path
from typing import List, Dict, Any, Tuple

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from reportlab.lib.pagesizes import landscape as rl_landscape, A4
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    PageBreak, Image as RLImage, KeepTogether,
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT, TA_CENTER
from reportlab.pdfgen import canvas as rl_canvas

from PIL import Image as PILImage

from app.core.paths import Paths

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _parse_slide_range(range_str: str, total: int) -> List[int]:
    """Parse '1-5, 7, 9-11' → list of 0-based indices within [0, total)."""
    if not range_str or range_str.strip().lower() in ("", "all"):
        return list(range(total))

    indices = set()
    for part in range_str.split(","):
        part = part.strip()
        if "-" in part:
            try:
                s, e = part.split("-", 1)
                for i in range(int(s.strip()), int(e.strip()) + 1):
                    if 1 <= i <= total:
                        indices.add(i - 1)
            except ValueError:
                pass
        else:
            try:
                i = int(part)
                if 1 <= i <= total:
                    indices.add(i - 1)
            except ValueError:
                pass
    return sorted(indices)


def _rgb_from_pptx(color) -> Tuple[float, float, float]:
    """Return an (r, g, b) tuple in 0-1 range from a pptx color object."""
    try:
        rgb: RGBColor = color.rgb
        return rgb.red / 255, rgb.green / 255, rgb.blue / 255
    except Exception:
        return 0, 0, 0


def _shape_text(shape) -> str:
    """Extract all text from a shape's text frame."""
    try:
        if shape.has_text_frame:
            return "\n".join(
                "".join(run.text for run in para.runs)
                for para in shape.text_frame.paragraphs
            )
    except Exception:
        pass
    return ""


def _slide_to_image(slide, width_px: int = 1280, height_px: int = 720) -> bytes | None:
    """
    Render a single slide as a PNG image using python-pptx's XML + PIL.
    Falls back to a simple white background if rendering fails.
    """
    try:
        # Try to render slide thumbnail via pptx if available (pptx >= 1.x)
        import tempfile, subprocess, sys

        # Use a simple approach: export shapes' images when present
        # For full fidelity rendering we would need LibreOffice or similar
        # Here we create a simple white image placeholder with slide dimensions
        img = PILImage.new("RGB", (width_px, height_px), color=(255, 255, 255))
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()
    except Exception as e:
        logger.warning(f"Slide image rendering failed: {e}")
        return None


# ---------------------------------------------------------------------------
# Main Service
# ---------------------------------------------------------------------------

class PowerPointToPdfService:
    def __init__(self):
        pass

    def _convert_single_file(
        self, input_path: Path, output_path: Path, config: Dict[str, Any]
    ) -> None:
        """
        Convert a single .pptx file to PDF using python-pptx + reportlab.
        Extracts slide text, titles, and embedded images for the PDF output.
        """
        prs = Presentation(str(input_path))
        total_slides = len(prs.slides)

        slide_range_str = config.get("slide_range", "all")
        include_hidden = config.get("include_hidden", False)
        layout_mode = config.get("layout_mode", "slides")
        quality = config.get("quality", "standard")

        selected_indices = _parse_slide_range(slide_range_str, total_slides)

        # Use landscape widescreen (16:9) as default page size for slides
        page_w = 10 * inch
        page_h = 7.5 * inch
        page_size = (page_w, page_h)

        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            "SlideTitle",
            parent=styles["Heading1"],
            fontSize=16,
            leading=20,
            spaceBefore=8,
            spaceAfter=4,
            textColor=colors.HexColor("#1a1a2e"),
        )
        body_style = ParagraphStyle(
            "SlideBody",
            parent=styles["Normal"],
            fontSize=11,
            leading=15,
            spaceAfter=3,
            textColor=colors.HexColor("#333333"),
        )
        notes_style = ParagraphStyle(
            "SlideNotes",
            parent=styles["Italic"],
            fontSize=8,
            leading=11,
            textColor=colors.grey,
        )
        slide_num_style = ParagraphStyle(
            "SlideNum",
            parent=styles["Normal"],
            fontSize=8,
            textColor=colors.HexColor("#999999"),
            alignment=TA_CENTER,
        )

        doc = SimpleDocTemplate(
            str(output_path),
            pagesize=page_size,
            leftMargin=0.5 * inch,
            rightMargin=0.5 * inch,
            topMargin=0.4 * inch,
            bottomMargin=0.4 * inch,
        )

        story = []

        for slide_idx in selected_indices:
            slide = prs.slides[slide_idx]

            # Skip hidden slides unless requested
            if not include_hidden:
                try:
                    if slide._element.get("show") == "0":
                        continue
                except Exception:
                    pass

            # Slide header line
            story.append(
                Paragraph(f"Slide {slide_idx + 1} of {total_slides}", slide_num_style)
            )

            # Collect shapes sorted top-to-bottom
            shapes = sorted(slide.shapes, key=lambda s: s.top if s.top is not None else 0)

            title_added = False
            for shape in shapes:
                # Embedded images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_blob = shape.image.blob
                        img_stream = io.BytesIO(img_blob)
                        pil_img = PILImage.open(img_stream)
                        pil_img.verify()  # check integrity
                        img_stream.seek(0)

                        # Scale image to fit within available width
                        max_w = page_w - inch
                        orig_w, orig_h = pil_img.size
                        scale = min(max_w / orig_w, 2.5 * inch / orig_h, 1.0)
                        draw_w = orig_w * scale
                        draw_h = orig_h * scale

                        story.append(
                            RLImage(img_stream, width=draw_w, height=draw_h)
                        )
                        story.append(Spacer(1, 4))
                    except Exception as img_err:
                        logger.debug(f"Could not embed image: {img_err}")
                    continue

                # Text shapes
                text = _shape_text(shape)
                if not text.strip():
                    continue

                # Detect title placeholder
                is_title = False
                try:
                    if shape.is_placeholder:
                        ph_type = shape.placeholder_format.type
                        if ph_type in (1, 3, 13):  # TITLE / CENTER_TITLE / SUBTITLE
                            is_title = True
                except Exception:
                    pass

                if is_title and not title_added:
                    story.append(Paragraph(text.strip(), title_style))
                    title_added = True
                else:
                    for line in text.split("\n"):
                        if line.strip():
                            safe_line = line.strip().replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                            story.append(Paragraph(safe_line, body_style))

            # Notes
            if layout_mode == "notes":
                try:
                    notes_frame = slide.notes_slide.notes_text_frame
                    notes_text = notes_frame.text.strip() if notes_frame else ""
                    if notes_text:
                        story.append(Spacer(1, 6))
                        story.append(Paragraph("Notes:", notes_style))
                        for n_line in notes_text.split("\n"):
                            if n_line.strip():
                                safe = n_line.strip().replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                                story.append(Paragraph(safe, notes_style))
                except Exception:
                    pass

            story.append(Spacer(1, 10))
            story.append(PageBreak())

        # Remove trailing PageBreak
        if story and isinstance(story[-1], PageBreak):
            story.pop()

        if not story:
            story = [Paragraph("No slides to display.", styles["Normal"])]

        doc.build(story)

    async def analyze(self, request_id: str, filename: str) -> Dict[str, Any]:
        """Extract slide count from a .pptx file using python-pptx."""
        upload_dir = Paths.request_upload(request_id)
        input_path = upload_dir / filename

        if not input_path.exists():
            raise ValueError("File not found.")

        try:
            prs = Presentation(str(input_path))
            slide_count = len(prs.slides)
            return {"filename": filename, "slide_count": slide_count}
        except Exception as e:
            logger.error(f"Failed to analyze PowerPoint file {filename}: {str(e)}")
            raise ValueError(f"Failed to read PowerPoint file: {str(e)}")

    async def process(
        self, request_id: str, filenames: List[str], config: Dict[str, Any]
    ) -> Dict[str, Any]:
        """Convert PowerPoint files to PDF (no Microsoft PowerPoint required)."""
        upload_dir = Paths.request_upload(request_id)
        output_dir = Paths.request_output(request_id)
        output_dir.mkdir(parents=True, exist_ok=True)

        results = []

        for filename in filenames:
            input_path = upload_dir / filename
            if not input_path.exists():
                results.append({
                    "original_filename": filename,
                    "status": "error",
                    "message": "File not found on server. Please re-upload.",
                })
                continue

            pdf_filename = config.get("output_filename")
            if not pdf_filename or len(filenames) > 1:
                pdf_filename = f"{Path(filename).stem}.pdf"
            if not pdf_filename.lower().endswith(".pdf"):
                pdf_filename += ".pdf"

            output_path = output_dir / pdf_filename

            try:
                self._convert_single_file(input_path, output_path, config)

                if not output_path.exists():
                    raise RuntimeError("PDF was not created.")

                results.append({
                    "original_filename": filename,
                    "pdf_filename": pdf_filename,
                    "status": "success",
                })
                logger.info(f"Successfully converted {filename} -> {pdf_filename}")

            except Exception as e:
                logger.error(f"Error converting {filename}: {str(e)}")
                results.append({
                    "original_filename": filename,
                    "status": "error",
                    "message": str(e),
                })

        return {
            "success": any(r.get("status") == "success" for r in results),
            "request_id": request_id,
            "results": results,
        }


powerpoint_to_pdf_service = PowerPointToPdfService()
