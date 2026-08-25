"""
PDF to PowerPoint conversion service.
Renders each PDF page as a high-quality image and inserts it into a PPTX slide.
"""

import logging
from pathlib import Path
from typing import Any, Dict
import io

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from app.core.paths import Paths

logger = logging.getLogger(__name__)

DPI = 150  # Render quality — higher = bigger file but better quality


class PDFToPowerPointService:
    async def process(
        self,
        request_id: str,
        filename: str,
        config: Dict[str, Any] = None,
    ) -> Dict[str, Any]:
        """Render each PDF page as an image and embed into a PPTX slide."""
        if config is None:
            config = {}

        upload_dir = Paths.request_upload(request_id)
        output_dir = Paths.request_output(request_id)
        output_dir.mkdir(parents=True, exist_ok=True)

        pdf_path = upload_dir / filename
        if not pdf_path.exists():
            raise ValueError(f"File not found: {filename}")

        out_name = f"{pdf_path.stem}.pptx"
        out_path = output_dir / out_name

        logger.info(f"Converting PDF to PowerPoint: {pdf_path} -> {out_path}")

        try:
            doc = fitz.open(str(pdf_path))
            prs = Presentation()

            # Use widescreen 16:9 slide layout
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            blank_layout = prs.slide_layouts[6]  # Blank layout

            zoom = DPI / 72.0  # PyMuPDF default is 72 DPI
            mat = fitz.Matrix(zoom, zoom)

            for page_num in range(len(doc)):
                page = doc[page_num]
                pix = page.get_pixmap(matrix=mat, alpha=False)

                # Convert pixmap to PNG bytes
                img_bytes = pix.tobytes("png")

                slide = prs.slides.add_slide(blank_layout)

                # Add image to fill the entire slide
                img_stream = io.BytesIO(img_bytes)
                slide.shapes.add_picture(
                    img_stream,
                    left=0,
                    top=0,
                    width=prs.slide_width,
                    height=prs.slide_height,
                )

            doc.close()
            prs.save(str(out_path))

        except Exception as e:
            logger.error(f"PDF to PowerPoint conversion failed: {e}", exc_info=True)
            raise ValueError(f"Failed to convert PDF to PowerPoint: {e}")

        if not out_path.exists():
            raise ValueError("Conversion succeeded but output file is missing.")

        return {
            "success": True,
            "request_id": request_id,
            "output_filename": out_name,
            "original_filename": filename,
        }


pdf_to_powerpoint_service = PDFToPowerPointService()
